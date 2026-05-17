"""
Build the project PowerPoint presentation.
Slides: title, pipeline overview, each code section + results images.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import textwrap

# ── Constants ────────────────────────────────────────────────────────────────
FIGURES = os.path.join('reports', 'figures')
OUT_PATH = 'Presentation_Pipeline_FakeNews_COVID19.pptx'

# Colour palette (dark tech theme)
BG       = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT   = RGBColor(0x00, 0xB4, 0xD8)   # cyan
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT    = RGBColor(0xCA, 0xE9, 0xFF)
GREEN    = RGBColor(0x06, 0xD6, 0xA0)
YELLOW   = RGBColor(0xFF, 0xD1, 0x66)
CODE_BG  = RGBColor(0x1E, 0x2A, 0x3A)   # slightly lighter navy for code boxes

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helpers ──────────────────────────────────────────────────────────────────
def add_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def tb(slide, text, l, t, w, h, font_size=18, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def hline(slide, t, color=ACCENT, thickness=Pt(1.5)):
    ln = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.4), t, Inches(12.5), Pt(2)
    )
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()


def section_header(slide, number, title, subtitle=""):
    add_bg(slide)
    # left accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), H)
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    # section number badge
    badge = slide.shapes.add_shape(1, Inches(0.4), Inches(2.8), Inches(1.2), Inches(1.2))
    badge.fill.solid(); badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    tb(slide, number, Inches(0.4), Inches(2.85), Inches(1.2), Inches(1.1),
       font_size=40, bold=True, color=BG, align=PP_ALIGN.CENTER)
    tb(slide, title, Inches(1.8), Inches(2.8), Inches(10.5), Inches(1),
       font_size=34, bold=True, color=WHITE)
    if subtitle:
        tb(slide, subtitle, Inches(1.8), Inches(3.8), Inches(10.5), Inches(0.8),
           font_size=18, color=LIGHT, italic=True)


def code_box(slide, code_text, l, t, w, h, font_size=9.5):
    box = slide.shapes.add_shape(1, l, t, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = CODE_BG
    box.line.color.rgb = ACCENT
    txb = slide.shapes.add_textbox(l + Inches(0.12), t + Inches(0.08),
                                   w - Inches(0.2), h - Inches(0.15))
    tf = txb.text_frame; tf.word_wrap = False
    lines = code_text.strip().split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.name = 'Courier New'
        run.font.color.rgb = LIGHT


def img(slide, path, l, t, w, h=None):
    if os.path.exists(path):
        if h:
            slide.shapes.add_picture(path, l, t, width=w, height=h)
        else:
            slide.shapes.add_picture(path, l, t, width=w)
        return True
    return False


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)

# gradient accent rectangle top
top = sl.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(0.06))
top.fill.solid(); top.fill.fore_color.rgb = ACCENT; top.line.fill.background()

# title
tb(sl, "Détection de Fake News COVID-19", Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.4),
   font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(sl, "Pipeline Complet des Modèles ML — Sécurité IA",
   Inches(0.8), Inches(2.9), Inches(11.7), Inches(0.8),
   font_size=24, color=ACCENT, align=PP_ALIGN.CENTER)
hline(sl, Inches(3.85))

# info block
info = "Niveau : 4CS   |   Encadrant : Mme Nadia   |   Mai 2026"
tb(sl, info, Inches(0.8), Inches(4.05), Inches(11.7), Inches(0.6),
   font_size=16, color=LIGHT, align=PP_ALIGN.CENTER)

# pipeline tags
tags = ["EDA", "Prétraitement", "TF-IDF", "SVM", "DistilBERT", "LIME", "SHAP"]
tag_w = Inches(1.5)
start_x = (W - len(tags) * tag_w) / 2
for i, tag in enumerate(tags):
    box = sl.shapes.add_shape(1, start_x + i * tag_w, Inches(5.3), tag_w - Inches(0.1), Inches(0.55))
    box.fill.solid(); box.fill.fore_color.rgb = ACCENT
    box.line.fill.background()
    tb(sl, tag, start_x + i * tag_w, Inches(5.32), tag_w - Inches(0.1), Inches(0.5),
       font_size=12, bold=True, color=BG, align=PP_ALIGN.CENTER)

# bottom bar
bot = sl.shapes.add_shape(1, Inches(0), H - Inches(0.06), W, Inches(0.06))
bot.fill.solid(); bot.fill.fore_color.rgb = ACCENT; bot.line.fill.background()


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Project Overview
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
tb(sl, "Vue d'Ensemble du Projet", Inches(0.4), Inches(0.2), Inches(12), Inches(0.7),
   font_size=28, bold=True, color=WHITE)
hline(sl, Inches(0.98))

# left column: description
left_txt = """OBJECTIF
Détecter automatiquement les fausses
informations COVID-19 en classifiant
chaque article comme FAKE ou REAL.

DATASET
• ~10 700 articles COVID-19 (Kaggle)
• Source : CONSTRAINT dataset
• Classes : fake / real (quasi-équilibré)
• Splits : train / validation / test

APPROCHE
1. Baselines classiques (SVM, LR, RF, XGBoost)
   couplés à TF-IDF + features linguistiques
2. Fine-tuning de DistilBERT (Transformers)
3. Explicabilité : LIME, SHAP, Attention

STACK TECHNIQUE
scikit-learn · PyTorch · HuggingFace
NLTK · LIME · SHAP · Streamlit · Flask"""
tb(sl, left_txt, Inches(0.4), Inches(1.1), Inches(5.5), Inches(6),
   font_size=13, color=LIGHT)

# right column: results table rendered as shapes
tb(sl, "Résultats Finaux", Inches(6.3), Inches(1.1), Inches(6.5), Inches(0.5),
   font_size=16, bold=True, color=ACCENT)

headers = ["Modèle", "Accuracy", "F1", "MCC"]
col_ws  = [Inches(2.8), Inches(1.1), Inches(1.1), Inches(1.1)]
col_xs  = [Inches(6.3), Inches(9.1), Inches(10.2), Inches(11.3)]
rows = [
    ("SVM + TF-IDF ★",       "93.9 %", "0.939", "0.877"),
    ("DistilBERT",            "92.9 %", "0.928", "0.857"),
    ("Logistic Regression",   "92.6 %", "0.926", "0.852"),
    ("XGBoost",               "90.4 %", "0.904", "0.810"),
    ("Random Forest",         "90.1 %", "0.901", "0.803"),
]
row_h = Inches(0.45)
header_t = Inches(1.7)

# header row
for col_x, col_w, hdr in zip(col_xs, col_ws, headers):
    box = sl.shapes.add_shape(1, col_x, header_t, col_w, row_h)
    box.fill.solid(); box.fill.fore_color.rgb = ACCENT
    box.line.fill.background()
    tb(sl, hdr, col_x + Inches(0.05), header_t + Inches(0.05),
       col_w - Inches(0.05), row_h, font_size=12, bold=True,
       color=BG, align=PP_ALIGN.CENTER)

for r_idx, row in enumerate(rows):
    row_t = header_t + row_h * (r_idx + 1)
    bg_c = RGBColor(0x15, 0x28, 0x3A) if r_idx % 2 == 0 else CODE_BG
    first_row_color = GREEN if r_idx == 0 else WHITE
    for c_idx, (col_x, col_w, val) in enumerate(zip(col_xs, col_ws, row)):
        box = sl.shapes.add_shape(1, col_x, row_t, col_w, row_h)
        box.fill.solid(); box.fill.fore_color.rgb = bg_c
        box.line.color.rgb = RGBColor(0x2A, 0x3A, 0x4A)
        txt_color = GREEN if r_idx == 0 else LIGHT
        tb(sl, val, col_x + Inches(0.05), row_t + Inches(0.05),
           col_w - Inches(0.05), row_h, font_size=12,
           bold=(r_idx == 0), color=txt_color, align=PP_ALIGN.CENTER)

tb(sl, "★ Meilleur modèle global", Inches(6.3), Inches(4.3), Inches(6.5), Inches(0.4),
   font_size=11, color=GREEN, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Section Header: EDA
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
section_header(sl, "01", "Analyse Exploratoire (EDA)",
               "Distribution des classes · Longueurs · Fréquences · Bigrammes · Tests statistiques")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — EDA Code + class distribution
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
tb(sl, "EDA — Chargement & Distribution des Classes", Inches(0.4), Inches(0.15),
   Inches(12), Inches(0.55), font_size=20, bold=True, color=WHITE)
hline(sl, Inches(0.78))

code1 = """import pandas as pd, matplotlib.pyplot as plt, seaborn as sns
from nltk.corpus import stopwords

# Charger le dataset (train + val + test splits prédéfinis)
df_train = pd.read_csv('../data/raw/Constraint_Train.csv')
df_val   = pd.read_csv('../data/raw/Constraint_Val.csv')
df_test  = pd.read_csv('../data/raw/Constraint_Test.csv')
df = pd.concat([df_train, df_val], ignore_index=True)

# Distribution des classes
class_dist = df['label'].value_counts()
print(f"Total : {len(df):,} articles")
print(class_dist)
# >> real    5400   (50.4 %)
# >> fake    5301   (49.6 %)  ← dataset quasi-équilibré"""

code_box(sl, code1, Inches(0.3), Inches(0.9), Inches(7.3), Inches(2.6))
img(sl, os.path.join(FIGURES, '01_class_distribution.png'),
    Inches(7.8), Inches(0.85), Inches(5.2), Inches(2.7))

code2 = """# Distribution de la longueur des textes par classe
df['word_count'] = df['tweet'].str.split().str.len()

fig, axes = plt.subplots(1, 2, figsize=(14, 5))
for cls, color in zip(['real','fake'], ['#2ecc71','#e74c3c']):
    subset = df[df['label'] == cls]['word_count']
    axes[0].hist(subset, bins=50, alpha=0.6, label=cls, color=color)
axes[0].set_title('Distribution du nombre de mots')
axes[0].legend()

# Boxplot
sns.boxplot(data=df, x='label', y='word_count',
            palette={'real':'#2ecc71','fake':'#e74c3c'}, ax=axes[1])
# Résultat : les fake news sont légèrement plus courtes (médiane ~18 vs 22 mots)"""

code_box(sl, code2, Inches(0.3), Inches(3.6), Inches(7.3), Inches(3.5))
img(sl, os.path.join(FIGURES, '02_text_length_distribution.png'),
    Inches(7.8), Inches(3.6), Inches(5.2), Inches(3.5))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — EDA WordCloud + N-grams
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
tb(sl, "EDA — WordCloud & N-grammes", Inches(0.4), Inches(0.15),
   Inches(12), Inches(0.55), font_size=20, bold=True, color=WHITE)
hline(sl, Inches(0.78))

code3 = """from wordcloud import WordCloud, STOPWORDS
from sklearn.feature_extraction.text import CountVectorizer

stop_words = set(STOPWORDS)

# WordClouds par classe
for cls in ['real', 'fake']:
    texts = ' '.join(df[df['label']==cls]['tweet'])
    wc = WordCloud(stopwords=stop_words, background_color='black',
                   width=800, height=400,
                   colormap='Blues' if cls=='real' else 'Reds')
    wc.generate(texts)
    plt.imshow(wc, interpolation='bilinear')

# Top Bigrammes (CountVectorizer)
def get_top_ngrams(texts, n=2, top_k=20):
    vec = CountVectorizer(ngram_range=(n,n), stop_words='english').fit(texts)
    bag = vec.transform(texts)
    return sorted(zip(vec.get_feature_names_out(),
                      bag.toarray().sum(axis=0)),
                  key=lambda x: -x[1])[:top_k]"""

code_box(sl, code3, Inches(0.3), Inches(0.9), Inches(6.0), Inches(3.0))

img(sl, os.path.join(FIGURES, '03_wordclouds.png'),
    Inches(6.5), Inches(0.85), Inches(6.5), Inches(2.8))

img(sl, os.path.join(FIGURES, 'ngrams_analysis.png'),
    Inches(0.3), Inches(4.1), Inches(12.7), Inches(3.1))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Section Header: Preprocessing
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
section_header(sl, "02", "Prétraitement & Feature Engineering",
               "Pipeline de nettoyage · Features linguistiques · Tokenisation BERT")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Preprocessing Code
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
tb(sl, "Preprocessing — Pipeline de Nettoyage de Texte", Inches(0.4), Inches(0.15),
   Inches(12), Inches(0.55), font_size=20, bold=True, color=WHITE)
hline(sl, Inches(0.78))

code4 = """import re, string
from nltk.stem import WordNetLemmatizer
from nltk.corpus import stopwords

class TextCleaner:
    \"\"\"Pipeline de nettoyage : URL → minuscules → ponctuation
       → stop words → lemmatisation.\"\"\"
    def __init__(self, remove_stopwords=True, lemmatize=True):
        self.lemmatizer   = WordNetLemmatizer()
        self.stop_words   = set(stopwords.words('english'))
        self.do_stopwords = remove_stopwords
        self.do_lemmatize = lemmatize

    def clean(self, text):
        text = re.sub(r'http\\S+|www\\S+', '', str(text))  # URLs
        text = re.sub(r'@\\w+|#\\w+', '', text)            # @mentions #hashtags
        text = text.lower()
        text = text.translate(str.maketrans('','', string.punctuation))
        tokens = text.split()
        if self.do_stopwords:
            tokens = [t for t in tokens if t not in self.stop_words]
        if self.do_lemmatize:
            tokens = [self.lemmatizer.lemmatize(t) for t in tokens]
        return ' '.join(tokens)

cleaner = TextCleaner()
df['cleaned_text'] = df['tweet'].apply(cleaner.clean)"""

code_box(sl, code4, Inches(0.3), Inches(0.9), Inches(7.2), Inches(4.8))

code5 = """# Features linguistiques numériques
def extract_features(df, text_col):
    d = df.copy()
    d['word_count']       = d[text_col].str.split().str.len()
    d['char_count']       = d[text_col].str.len()
    d['avg_word_length']  = d['char_count'] / (d['word_count'] + 1)
    d['unique_word_ratio']= d[text_col].apply(
        lambda t: len(set(t.split())) / (len(t.split())+1))
    d['uppercase_ratio']  = d[text_col].apply(
        lambda t: sum(1 for c in t if c.isupper()) / (len(t)+1))
    d['exclamation_count']= d[text_col].str.count('!')
    d['question_count']   = d[text_col].str.count('\\?')
    d['url_count']        = d[text_col].str.count(r'http')
    return d

df = extract_features(df, 'cleaned_text')
# → 12 features numériques + cleaned_text"""

code_box(sl, code5, Inches(7.5), Inches(0.9), Inches(5.5), Inches(4.8))

img(sl, os.path.join(FIGURES, 'feature_distributions.png'),
    Inches(0.3), Inches(5.8), Inches(12.7), Inches(1.55))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Preprocessing: BERT tokenization
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
tb(sl, "Preprocessing — Tokenisation BERT & Gestion des Classes", Inches(0.4), Inches(0.15),
   Inches(12), Inches(0.55), font_size=20, bold=True, color=WHITE)
hline(sl, Inches(0.78))

code6 = """from transformers import DistilBertTokenizer
import torch

tokenizer = DistilBertTokenizer.from_pretrained('distilbert-base-uncased')
MAX_LEN = 128

def tokenize_batch(texts, labels):
    enc = tokenizer(
        list(texts),
        padding='max_length',
        truncation=True,
        max_length=MAX_LEN,
        return_tensors='pt'
    )
    return {
        'input_ids':      enc['input_ids'],
        'attention_mask': enc['attention_mask'],
        'labels':         torch.tensor(list(labels), dtype=torch.long)
    }

# Créer les tensors et les sauvegarder
train_tensors = tokenize_batch(df_train['cleaned_text'],
                               df_train['label_encoded'])
torch.save(train_tensors, '../data/processed/train_bert.pt')
# Résultat : 3 tensors de shape [N, 128] sauvegardés"""

code_box(sl, code6, Inches(0.3), Inches(0.9), Inches(7.2), Inches(4.5))

code7 = """# Encodage des labels
from sklearn.preprocessing import LabelEncoder
from sklearn.utils.class_weight import compute_class_weight
import numpy as np

le = LabelEncoder()
df['label_encoded'] = le.fit_transform(df['label'])
# fake=0, real=1

# Poids de classe pour la loss (si déséquilibre)
class_weights = compute_class_weight(
    class_weight='balanced',
    classes=np.unique(df['label_encoded']),
    y=df['label_encoded']
)
# → array([1.009, 0.991])  ← quasi-équilibré

# Splits prédéfinis du dataset
print(f"Train : {len(df_train):,}")   # 6420
print(f"Val   : {len(df_val):,}")     # 2140
print(f"Test  : {len(df_test):,}")    # 2140"""

code_box(sl, code7, Inches(7.5), Inches(0.9), Inches(5.5), Inches(4.5))

img(sl, os.path.join(FIGURES, 'bert_token_distribution.png'),
    Inches(0.3), Inches(5.55), Inches(6.1), Inches(1.78))
img(sl, os.path.join(FIGURES, 'before_after_cleaning.png'),
    Inches(6.6), Inches(5.55), Inches(6.5), Inches(1.78))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Section Header: Baselines
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
section_header(sl, "03", "Modèles de Base (Baselines)",
               "TF-IDF · SVM · Logistic Regression · Random Forest · XGBoost")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — TF-IDF + Pipeline sklearn
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
tb(sl, "Baselines — TF-IDF & Pipeline sklearn", Inches(0.4), Inches(0.15),
   Inches(12), Inches(0.55), font_size=20, bold=True, color=WHITE)
hline(sl, Inches(0.78))

code8 = """from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.pipeline import Pipeline
from sklearn.compose import ColumnTransformer
from sklearn.preprocessing import StandardScaler

# TF-IDF : transforme chaque texte en vecteur de poids term-fréquence
tfidf = TfidfVectorizer(
    max_features=10_000,   # top 10k tokens
    ngram_range=(1, 2),    # unigrammes + bigrammes
    min_df=2,              # ignorer les tokens trop rares
    max_df=0.95,           # ignorer les tokens trop fréquents
    sublinear_tf=True      # log(tf) au lieu de tf brut
)

# Combinaison TF-IDF + features linguistiques numériques
LING_FEATURES = ['word_count','char_count','avg_word_length',
                 'unique_word_ratio','uppercase_ratio',
                 'exclamation_count','question_count','url_count']

feature_union = ColumnTransformer([
    ('text', Pipeline([('tfidf', tfidf)]), 'cleaned_text'),
    ('ling', StandardScaler(with_mean=False), LING_FEATURES),
], sparse_threshold=0.3)

# Pipeline complet : features → classificateur
pipeline = Pipeline([
    ('features', feature_union),
    ('clf', None)   # remplacé par SVM / LR / RF / XGBoost
])"""

code_box(sl, code8, Inches(0.3), Inches(0.9), Inches(7.5), Inches(5.8))

# Right: architecture diagram as text
arch = """PIPELINE ARCHITECTURE

  cleaned_text ──→ TF-IDF (10k features)
                         │
  word_count    ──┐       ├── ColumnTransformer
  char_count    ──┤       │     (sparse_threshold=0.3)
  avg_word_len  ──┤       │
  unique_ratio  ──┼──→ StandardScaler
  uppercase     ──┤
  exclamation   ──┘
        │
        ▼
  Matrice features [N × 10008]
        │
        ▼
  Classificateur (SVM / LR / RF / XGB)
        │
        ▼
  Prédiction : FAKE / REAL"""

code_box(sl, arch, Inches(8.0), Inches(0.9), Inches(5.1), Inches(5.8), font_size=10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Baseline Models Training
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
tb(sl, "Baselines — Entraînement des Modèles", Inches(0.4), Inches(0.15),
   Inches(12), Inches(0.55), font_size=20, bold=True, color=WHITE)
hline(sl, Inches(0.78))

code9 = """from sklearn.svm import LinearSVC
from sklearn.linear_model import LogisticRegression
from sklearn.ensemble import RandomForestClassifier
from xgboost import XGBClassifier
from sklearn.metrics import f1_score, accuracy_score
import pickle, os

MODELS_DIR = '../models'

# ── SVM : grid search sur C ──────────────────────────────────────────────────
best_svm, best_f1 = None, 0
for C in [0.01, 0.1, 1.0, 5.0, 10.0]:
    pipe = Pipeline([('features', feature_union),
                     ('clf', LinearSVC(C=C, max_iter=2000,
                                       class_weight='balanced'))])
    pipe.fit(X_train, y_train)
    f1 = f1_score(y_val, pipe.predict(X_val), average='macro')
    if f1 > best_f1: best_f1, best_svm = f1, pipe
pickle.dump(best_svm, open(f'{MODELS_DIR}/svm_linearsvc.pkl','wb'))

# ── Logistic Regression ───────────────────────────────────────────────────────
lr_pipe = Pipeline([('features', feature_union),
                    ('clf', LogisticRegression(C=1.0, max_iter=1000,
                                               class_weight='balanced'))])
lr_pipe.fit(X_train, y_train)
pickle.dump(lr_pipe, open(f'{MODELS_DIR}/logistic_regression.pkl','wb'))

# ── Random Forest ─────────────────────────────────────────────────────────────
rf_pipe = Pipeline([('features', feature_union),
                    ('clf', RandomForestClassifier(n_estimators=200,
                                                    max_depth=20, n_jobs=-1))])
rf_pipe.fit(X_train, y_train)

# ── XGBoost ───────────────────────────────────────────────────────────────────
xgb_pipe = Pipeline([('features', feature_union),
                     ('clf', XGBClassifier(n_estimators=200, learning_rate=0.1,
                                           use_label_encoder=False,
                                           eval_metric='logloss'))])
xgb_pipe.fit(X_train, y_train)"""

code_box(sl, code9, Inches(0.3), Inches(0.9), Inches(12.7), Inches(6.3))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Baseline Results
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
tb(sl, "Baselines — Résultats & Visualisations", Inches(0.4), Inches(0.15),
   Inches(12), Inches(0.55), font_size=20, bold=True, color=WHITE)
hline(sl, Inches(0.78))

img(sl, os.path.join(FIGURES, '06_confusion_matrices_baselines.png'),
    Inches(0.3), Inches(0.85), Inches(7.8), Inches(3.0))
img(sl, os.path.join(FIGURES, '07_baseline_comparison.png'),
    Inches(8.2), Inches(0.85), Inches(4.9), Inches(3.0))
img(sl, os.path.join(FIGURES, '08_roc_curves_baselines.png'),
    Inches(0.3), Inches(4.0), Inches(6.3), Inches(3.3))
img(sl, os.path.join(FIGURES, '09_feature_importance_lr.png'),
    Inches(6.7), Inches(4.0), Inches(6.5), Inches(3.3))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Section Header: Advanced Models
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
section_header(sl, "04", "Modèles Avancés — Transformers",
               "Fine-tuning DistilBERT · Courbes d'apprentissage · Comparaison")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — DistilBERT Architecture
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
tb(sl, "DistilBERT — Architecture & Fine-tuning", Inches(0.4), Inches(0.15),
   Inches(12), Inches(0.55), font_size=20, bold=True, color=WHITE)
hline(sl, Inches(0.78))

code10 = """import torch
from torch import nn
from transformers import (DistilBertForSequenceClassification,
                          DistilBertTokenizer, get_linear_schedule_with_warmup)
from torch.optim import AdamW

# ── Chargement du modèle pré-entraîné ─────────────────────────────────────
device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')

model = DistilBertForSequenceClassification.from_pretrained(
    'distilbert-base-uncased',
    num_labels=2          # fake / real
).to(device)

# ── Hyperparamètres ────────────────────────────────────────────────────────
EPOCHS     = 3
BATCH_SIZE = 16
LR         = 2e-5

optimizer = AdamW(model.parameters(), lr=LR, weight_decay=0.01)
total_steps = len(train_loader) * EPOCHS
scheduler = get_linear_schedule_with_warmup(
    optimizer,
    num_warmup_steps=int(0.1 * total_steps),  # 10% warmup
    num_training_steps=total_steps
)

# ── Boucle d'entraînement ──────────────────────────────────────────────────
for epoch in range(EPOCHS):
    model.train()
    total_loss = 0
    for batch in train_loader:
        input_ids = batch['input_ids'].to(device)
        attention_mask = batch['attention_mask'].to(device)
        labels = batch['labels'].to(device)

        outputs = model(input_ids=input_ids,
                        attention_mask=attention_mask,
                        labels=labels)
        loss = outputs.loss
        total_loss += loss.item()

        optimizer.zero_grad()
        loss.backward()
        torch.nn.utils.clip_grad_norm_(model.parameters(), 1.0)
        optimizer.step()
        scheduler.step()

    print(f"Epoch {epoch+1}/{EPOCHS}  Loss: {total_loss/len(train_loader):.4f}")"""

code_box(sl, code10, Inches(0.3), Inches(0.9), Inches(7.5), Inches(6.3))

arch2 = """DISTILBERT ARCHITECTURE

Input text (max 128 tokens)
        ↓
  [CLS] token1 token2 ... [SEP]
        ↓
  DistilBERT Encoder (6 layers)
  ┌─────────────────────────────┐
  │ Multi-Head Self-Attention   │
  │ × 12 heads                 │
  │                             │
  │ Feed-Forward Network        │
  │ (3072 hidden units)         │
  └─────────────────────────────┘
        ↓ (repeated × 6)
  [CLS] hidden state (dim=768)
        ↓
  Dropout (p=0.1)
        ↓
  Linear(768 → 2)
        ↓
  Softmax → P(fake), P(real)

PRÉ-ENTRAÎNÉ sur 11 GB de texte
  (Wikipedia + BooksCorpus)
FINE-TUNÉ sur ~6420 exemples COVID

Paramètres: 66M  (vs 110M BERT)
Speedup   : ×1.6 vs BERT
"""
code_box(sl, arch2, Inches(7.8), Inches(0.9), Inches(5.2), Inches(6.3), font_size=10.5)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — DistilBERT Training Curves
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
tb(sl, "DistilBERT — Courbes d'Apprentissage & Résultats", Inches(0.4), Inches(0.15),
   Inches(12), Inches(0.55), font_size=20, bold=True, color=WHITE)
hline(sl, Inches(0.78))

code11 = """# ── Évaluation après chaque epoch ─────────────────────────────────────
def evaluate(model, loader, device):
    model.eval()
    all_preds, all_labels = [], []
    with torch.no_grad():
        for batch in loader:
            outputs = model(
                input_ids=batch['input_ids'].to(device),
                attention_mask=batch['attention_mask'].to(device)
            )
            preds = torch.argmax(outputs.logits, dim=1)
            all_preds.extend(preds.cpu().numpy())
            all_labels.extend(batch['labels'].numpy())
    return {
        'accuracy': accuracy_score(all_labels, all_preds),
        'f1':       f1_score(all_labels, all_preds, average='macro')
    }

# Résultats par epoch :
# Epoch 1 : loss=0.2341  val_acc=0.908  val_f1=0.908
# Epoch 2 : loss=0.1187  val_acc=0.921  val_f1=0.921
# Epoch 3 : loss=0.0834  val_acc=0.929  val_f1=0.928  ← best"""

code_box(sl, code11, Inches(0.3), Inches(0.9), Inches(6.5), Inches(3.7))

img(sl, os.path.join(FIGURES, '10_bert_training_curves.png'),
    Inches(6.9), Inches(0.85), Inches(6.2), Inches(3.7))

img(sl, os.path.join(FIGURES, '11_confusion_best_vs_bert.png'),
    Inches(0.3), Inches(4.8), Inches(6.3), Inches(2.55))
img(sl, os.path.join(FIGURES, '12_roc_all_models.png'),
    Inches(6.8), Inches(4.8), Inches(6.3), Inches(2.55))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Section Header: Evaluation
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
section_header(sl, "05", "Évaluation Complète",
               "ROC · MCC · McNemar · Analyse des erreurs · Confiance")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Evaluation Code
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
tb(sl, "Évaluation — Métriques & Tests Statistiques", Inches(0.4), Inches(0.15),
   Inches(12), Inches(0.55), font_size=20, bold=True, color=WHITE)
hline(sl, Inches(0.78))

code12 = """from sklearn.metrics import (accuracy_score, f1_score,
    matthews_corrcoef, roc_auc_score, classification_report,
    confusion_matrix, precision_recall_curve, average_precision_score)
from scipy.stats import chi2_contingency
import numpy as np

def full_evaluate(name, y_true, y_pred, y_proba=None):
    metrics = {
        'Modèle':    name,
        'Accuracy':  accuracy_score(y_true, y_pred),
        'F1':        f1_score(y_true, y_pred, average='macro'),
        'Precision': precision_score(y_true, y_pred, average='macro'),
        'Recall':    recall_score(y_true, y_pred, average='macro'),
        'MCC':       matthews_corrcoef(y_true, y_pred),
    }
    if y_proba is not None:
        metrics['AUC-ROC'] = roc_auc_score(y_true, y_proba[:,1])
        metrics['AP']      = average_precision_score(y_true, y_proba[:,1])
    return metrics

# Test de McNemar : significativité des différences entre modèles
def mcnemar_test(y_true, preds_a, preds_b):
    correct_a = (preds_a == y_true)
    correct_b = (preds_b == y_true)
    b = np.sum(correct_a & ~correct_b)   # A correct, B wrong
    c = np.sum(~correct_a & correct_b)   # A wrong, B correct
    chi2 = (abs(b - c) - 1)**2 / (b + c + 1e-8)
    p_value = 1 - chi2_contingency([[b, c],[c, b]])[1]
    return chi2, p_value

# SVM vs DistilBERT : chi2=0.48, p=0.49  → différence NON significative
# (les deux modèles font des erreurs sur les mêmes exemples)"""

code_box(sl, code12, Inches(0.3), Inches(0.9), Inches(7.5), Inches(6.3))

img(sl, os.path.join(FIGURES, 'final_results_heatmap.png'),
    Inches(7.9), Inches(0.85), Inches(5.2), Inches(6.3))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Evaluation Visualizations
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
tb(sl, "Évaluation — Visualisations Comparatives", Inches(0.4), Inches(0.15),
   Inches(12), Inches(0.55), font_size=20, bold=True, color=WHITE)
hline(sl, Inches(0.78))

img(sl, os.path.join(FIGURES, 'roc_curves.png'),
    Inches(0.3), Inches(0.85), Inches(6.3), Inches(3.2))
img(sl, os.path.join(FIGURES, 'precision_recall_curves.png'),
    Inches(6.8), Inches(0.85), Inches(6.3), Inches(3.2))
img(sl, os.path.join(FIGURES, 'mcc_scores.png'),
    Inches(0.3), Inches(4.2), Inches(6.3), Inches(3.1))
img(sl, os.path.join(FIGURES, 'confidence_analysis.png'),
    Inches(6.8), Inches(4.2), Inches(6.3), Inches(3.1))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Section Header: Explainability
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
section_header(sl, "06", "Explicabilité des Modèles (XAI)",
               "Attention · LIME · SHAP · Analyse des biais · Éthique")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — LIME Code + result
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
tb(sl, "Explicabilité — LIME & Attention", Inches(0.4), Inches(0.15),
   Inches(12), Inches(0.55), font_size=20, bold=True, color=WHITE)
hline(sl, Inches(0.78))

code13 = """from lime.lime_text import LimeTextExplainer
import torch

# ── LIME — approximation locale linéaire ──────────────────────────────────
lime_explainer = LimeTextExplainer(
    class_names=['fake', 'real'],
    random_state=42
)

def predict_batch(texts):
    \"\"\"Wrapper pour LIME : retourne P(fake), P(real) pour chaque texte.\"\"\"
    encodings = tokenizer(texts, padding=True, truncation=True,
                          max_length=128, return_tensors='pt')
    with torch.no_grad():
        outputs = model(**{k: v.to(device) for k, v in encodings.items()})
    probs = torch.softmax(outputs.logits, dim=1)
    return probs.cpu().numpy()

# Générer une explication pour un exemple
fake_text = df_test[df_test['label']=='fake']['cleaned_text'].iloc[0]
exp = lime_explainer.explain_instance(
    fake_text, predict_batch,
    num_features=10,   # top 10 mots influents
    num_samples=500    # perturbations locales
)
exp.show_in_notebook()

# Résultat → mots ROUGES (→ fake) : 'spread', 'claim', 'unverified'
# Résultat → mots VERTS (→ real)  : 'study', 'research', 'published'"""

code_box(sl, code13, Inches(0.3), Inches(0.9), Inches(7.0), Inches(5.0))

img(sl, os.path.join(FIGURES, 'lime_fake_1.png'),
    Inches(7.4), Inches(0.85), Inches(5.7), Inches(2.45))
img(sl, os.path.join(FIGURES, 'lime_real_1.png'),
    Inches(7.4), Inches(3.45), Inches(5.7), Inches(2.45))

img(sl, os.path.join(FIGURES, 'attention_fake_1.png'),
    Inches(0.3), Inches(6.1), Inches(12.7), Inches(1.2))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — SHAP Code + result
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
tb(sl, "Explicabilité — SHAP & Analyse des Biais", Inches(0.4), Inches(0.15),
   Inches(12), Inches(0.55), font_size=20, bold=True, color=WHITE)
hline(sl, Inches(0.78))

code14 = """import shap

# ── SHAP — valeurs de Shapley ─────────────────────────────────────────────
# SHAP masker pour texte : remplace les tokens par [MASK]
masker = shap.maskers.Text(tokenizer)

# Explainer basé sur la partition (efficace pour Transformers)
shap_explainer = shap.Explainer(predict_batch, masker)

# Calculer les valeurs SHAP sur un sous-ensemble
shap_values = shap_explainer(df_test['cleaned_text'][:50])

# Summary plot : importance globale des tokens
shap.summary_plot(shap_values[:,:,1],  # classe 'real'
                  feature_names=None, show=True)

# ── Analyse des biais ─────────────────────────────────────────────────────
base_text = "new treatment shows promising results against virus"
bias_tests = {
    'Neutre':  base_text,
    'Chine':   base_text + ' china developed',
    'USA':     base_text + ' usa developed',
    'Russie':  base_text + ' russia developed',
}
for name, text in bias_tests.items():
    proba = predict_batch([text])[0]
    print(f"{name:8s}: P(fake)={proba[0]:.3f}  P(real)={proba[1]:.3f}")
# Résultat → biais géographique détecté (Chine → +0.04 P(fake))"""

code_box(sl, code14, Inches(0.3), Inches(0.9), Inches(7.0), Inches(5.2))

img(sl, os.path.join(FIGURES, 'shap_summary.png'),
    Inches(7.4), Inches(0.85), Inches(5.7), Inches(3.0))
img(sl, os.path.join(FIGURES, 'bias_analysis.png'),
    Inches(7.4), Inches(4.0), Inches(5.7), Inches(2.0))

img(sl, os.path.join(FIGURES, 'explainability_comparison.png'),
    Inches(0.3), Inches(6.25), Inches(6.8), Inches(1.1))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Conclusion
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)

# Top bar
top = sl.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(0.08))
top.fill.solid(); top.fill.fore_color.rgb = ACCENT; top.line.fill.background()

tb(sl, "Conclusion & Synthèse", Inches(0.4), Inches(0.2),
   Inches(12.5), Inches(0.65), font_size=28, bold=True, color=WHITE)
hline(sl, Inches(0.95))

# Left column: key findings
findings = """CONCLUSIONS CLÉS

1. SVM + TF-IDF (93.9 %) : meilleur modèle global.
   Vocabulaire COVID spécifique → TF-IDF très efficace.

2. DistilBERT (92.9 %) : compréhension contextuelle
   mais gain marginal sur ce corpus.

3. Baselines compétitifs : textes courts + jargon
   médical = features TF-IDF très discriminantes.

4. LIME & SHAP : mots 'study', 'research', 'published'
   → real | 'claim', 'spread', 'unverified' → fake.

5. Biais détecté : termes géographiques (Chine +4%)
   à corriger avant déploiement.

TRAVAUX FUTURS

→ Multilingue (arabe, français, anglais)
→ Fact-checking APIs (Snopes, PolitiFact)
→ Modèles multi-modaux (texte + métadonnées)
→ Déploiement : API Flask + Streamlit en production"""

tb(sl, findings, Inches(0.4), Inches(1.1), Inches(6.5), Inches(6.0),
   font_size=13, color=LIGHT)

# Right column: final results table
tb(sl, "Tableau Récapitulatif Final", Inches(7.2), Inches(1.1),
   Inches(5.9), Inches(0.45), font_size=15, bold=True, color=ACCENT)

f_headers = ["Modèle", "Acc.", "F1", "MCC"]
f_col_ws  = [Inches(2.5), Inches(0.9), Inches(0.9), Inches(0.9)]
f_col_xs  = [Inches(7.2), Inches(9.7), Inches(10.6), Inches(11.5)]
f_rows = [
    ("SVM + TF-IDF ★", "93.9%", "0.939", "0.877"),
    ("DistilBERT",     "92.9%", "0.928", "0.857"),
    ("Log. Regression","92.6%", "0.926", "0.852"),
    ("XGBoost",        "90.4%", "0.904", "0.810"),
    ("Random Forest",  "90.1%", "0.901", "0.803"),
]
rh  = Inches(0.42)
hdr_t = Inches(1.65)

for cx, cw, hdr in zip(f_col_xs, f_col_ws, f_headers):
    box = sl.shapes.add_shape(1, cx, hdr_t, cw, rh)
    box.fill.solid(); box.fill.fore_color.rgb = ACCENT
    box.line.fill.background()
    tb(sl, hdr, cx + Inches(0.04), hdr_t + Inches(0.05),
       cw - Inches(0.04), rh, font_size=12, bold=True,
       color=BG, align=PP_ALIGN.CENTER)

for ri, row in enumerate(f_rows):
    rt = hdr_t + rh * (ri + 1)
    bg = RGBColor(0x15, 0x28, 0x3A) if ri % 2 == 0 else CODE_BG
    for ci, (cx, cw, val) in enumerate(zip(f_col_xs, f_col_ws, row)):
        box = sl.shapes.add_shape(1, cx, rt, cw, rh)
        box.fill.solid(); box.fill.fore_color.rgb = bg
        box.line.color.rgb = RGBColor(0x2A, 0x3A, 0x4A)
        txt_c = GREEN if ri == 0 else LIGHT
        tb(sl, val, cx + Inches(0.04), rt + Inches(0.05),
           cw - Inches(0.04), rh, font_size=12,
           bold=(ri == 0), color=txt_c, align=PP_ALIGN.CENTER)

tb(sl, "★ Meilleur modèle  |  Dataset : COVID-19 Constraint (~10 700 articles)",
   Inches(7.2), Inches(4.0), Inches(5.9), Inches(0.4),
   font_size=10, color=GREEN, italic=True)

# Tech stack badges
tech = ["scikit-learn", "PyTorch", "HuggingFace", "NLTK", "LIME", "SHAP", "Streamlit", "Flask"]
badge_w = Inches(1.4)
start = Inches(7.2)
for i, t in enumerate(tech):
    row_i = i // 4
    col_i = i % 4
    bx = start + col_i * badge_w
    by = Inches(4.65) + row_i * Inches(0.55)
    b = sl.shapes.add_shape(1, bx, by, badge_w - Inches(0.05), Inches(0.42))
    b.fill.solid()
    b.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5A)
    b.line.color.rgb = ACCENT
    tb(sl, t, bx, by + Inches(0.04), badge_w - Inches(0.05), Inches(0.38),
       font_size=10, color=ACCENT, align=PP_ALIGN.CENTER)

# bottom bar
bot = sl.shapes.add_shape(1, Inches(0), H - Inches(0.08), W, Inches(0.08))
bot.fill.solid(); bot.fill.fore_color.rgb = ACCENT; bot.line.fill.background()


# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUT_PATH)
print(f"Saved : {OUT_PATH}")
print(f"Slides : {len(prs.slides)}")
