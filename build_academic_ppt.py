"""
Academic PowerPoint — NO code, scientific article structure.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

FIGURES = os.path.join('reports', 'figures')
OUT = 'Presentation_Academique_FakeNews_COVID19.pptx'

BG    = RGBColor(0x0D, 0x1B, 0x2A)
ACC   = RGBColor(0x00, 0xB4, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xCA, 0xE9, 0xFF)
GREEN = RGBColor(0x06, 0xD6, 0xA0)
YELL  = RGBColor(0xFF, 0xD1, 0x66)
DARK  = RGBColor(0x15, 0x28, 0x3A)
RED   = RGBColor(0xFF, 0x6B, 0x6B)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── Helpers ──────────────────────────────────────────────────────────────────
def bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

def tb(slide, text, l, t, w, h, fs=14, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text = text
    r.font.size   = Pt(fs)
    r.font.bold   = bold
    r.font.color.rgb = color
    r.font.italic = italic
    return txb

def rect(slide, l, t, w, h, fill_rgb, line_rgb=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        s.line.color.rgb = line_rgb
    else:
        s.line.fill.background()
    return s

def hline(slide, y, color=ACC):
    r = rect(slide, Inches(0.35), y, Inches(12.6), Pt(2), color)

def img(slide, path, l, t, w, h=None):
    if os.path.exists(path):
        if h:
            slide.shapes.add_picture(path, l, t, width=w, height=h)
        else:
            slide.shapes.add_picture(path, l, t, width=w)

def top_bar(slide):
    r = rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACC)
def bot_bar(slide):
    r = rect(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), ACC)

def slide_title(slide, title, subtitle=""):
    top_bar(slide)
    rect(slide, Inches(0), Inches(0), Inches(0.12), H, ACC)
    tb(slide, title, Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.65),
       fs=22, bold=True, color=WHITE)
    hline(slide, Inches(0.82))
    if subtitle:
        tb(slide, subtitle, Inches(0.3), Inches(0.85), Inches(12.7), Inches(0.4),
           fs=12, color=LIGHT, italic=True)
    bot_bar(slide)

def section_divider(slide, number, title, subtitle=""):
    bg(slide)
    rect(slide, Inches(0), Inches(0), Inches(0.15), H, ACC)
    # large number
    r = rect(slide, Inches(0.5), Inches(2.5), Inches(1.4), Inches(1.4), ACC)
    r.line.fill.background()
    tb(slide, number, Inches(0.5), Inches(2.55), Inches(1.4), Inches(1.3),
       fs=44, bold=True, color=BG, align=PP_ALIGN.CENTER)
    tb(slide, title, Inches(2.2), Inches(2.5), Inches(10.5), Inches(1.1),
       fs=36, bold=True, color=WHITE)
    if subtitle:
        tb(slide, subtitle, Inches(2.2), Inches(3.65), Inches(10.5), Inches(0.7),
           fs=16, color=LIGHT, italic=True)

def bullet_box(slide, items, l, t, w, h, fs=12.5, title=None, title_color=ACC):
    if title:
        tb(slide, title, l, t, w, Inches(0.4), fs=13, bold=True, color=title_color)
        t += Inches(0.42)
        h -= Inches(0.42)
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run()
        if isinstance(item, tuple):
            r.text = f"• {item[0]}"
            r.font.bold = True
            r.font.color.rgb = ACC
            r.font.size = Pt(fs)
            r2 = p.add_run()
            r2.text = f"  {item[1]}"
            r2.font.size = Pt(fs - 0.5)
            r2.font.color.rgb = LIGHT
        else:
            r.text = f"• {item}"
            r.font.size = Pt(fs)
            r.font.color.rgb = LIGHT


def results_table(slide, l, t):
    headers = ["Modèle", "Accuracy", "F1", "MCC", "AUC-ROC"]
    rows = [
        ("SVM + TF-IDF ★", "93,9 %", "0,939", "0,877", "—"),
        ("DistilBERT",      "92,9 %", "0,928", "0,857", "0,976"),
        ("Log. Regression", "92,6 %", "0,926", "0,852", "0,976"),
        ("XGBoost",         "90,4 %", "0,904", "0,810", "0,968"),
        ("Random Forest",   "90,1 %", "0,901", "0,803", "0,969"),
    ]
    col_ws = [Inches(2.6), Inches(1.0), Inches(0.85), Inches(0.85), Inches(0.95)]
    col_xs = [l]
    for cw in col_ws[:-1]:
        col_xs.append(col_xs[-1] + cw)
    rh = Inches(0.4)

    # header
    for cx, cw, h in zip(col_xs, col_ws, headers):
        rect(slide, cx, t, cw, rh, ACC).line.fill.background()
        tb(slide, h, cx + Inches(0.04), t + Inches(0.04),
           cw - Inches(0.06), rh, fs=11, bold=True, color=BG, align=PP_ALIGN.CENTER)

    for ri, row in enumerate(rows):
        ry = t + rh * (ri + 1)
        bg_c = DARK if ri % 2 == 0 else RGBColor(0x1E, 0x2A, 0x3A)
        txt_c = GREEN if ri == 0 else LIGHT
        for ci, (cx, cw, val) in enumerate(zip(col_xs, col_ws, row)):
            b = rect(slide, cx, ry, cw, rh, bg_c, RGBColor(0x2A,0x3A,0x4A))
            tb(slide, val, cx + Inches(0.03), ry + Inches(0.04),
               cw - Inches(0.04), rh, fs=11,
               bold=(ri == 0), color=txt_c, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITRE
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
top_bar(sl)
bot_bar(sl)

# Main title
tb(sl,
   "Détection Automatique de Fausses Informations COVID-19\n"
   "par Apprentissage Automatique",
   Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.6),
   fs=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

hline(sl, Inches(3.0))

tb(sl,
   "Comparaison de modèles classiques et Transformers\navec analyse d'explicabilité (XAI)",
   Inches(0.8), Inches(3.1), Inches(11.7), Inches(0.9),
   fs=20, color=LIGHT, align=PP_ALIGN.CENTER, italic=True)

tb(sl,
   "Youssef Benmabrouk  |  4CS  |  Encadrant : Mme Nadia  |  Mai 2026",
   Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.5),
   fs=14, color=ACC, align=PP_ALIGN.CENTER)

# Bottom keywords
kws = ["Infodémie COVID-19", "NLP", "SVM", "DistilBERT", "TF-IDF", "LIME", "SHAP", "XAI"]
kw_w = Inches(1.55)
start = (W - len(kws) * kw_w) / 2
for i, kw in enumerate(kws):
    b = rect(sl, start + i * kw_w, Inches(5.6), kw_w - Inches(0.08), Inches(0.52), ACC)
    b.line.fill.background()
    tb(sl, kw, start + i * kw_w, Inches(5.62), kw_w - Inches(0.08), Inches(0.48),
       fs=11, bold=True, color=BG, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PLAN
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "Plan de la Présentation")

sections = [
    ("01", "Introduction", "Contexte · Infodémie · Problématique · Contributions"),
    ("02", "État de l'art", "Fake news detection 2017–2021 · BERT · Constraint"),
    ("03", "Dataset & Prétraitement", "Constraint COVID-19 · Nettoyage · TF-IDF · Features"),
    ("04", "Architectures des modèles", "SVM · LR · RF · XGBoost · DistilBERT"),
    ("05", "Résultats & Discussion", "Tableau comparatif · ROC · Matrices · Features"),
    ("06", "Explicabilité (XAI)", "LIME · Attention · Biais · Éthique"),
    ("07", "Conclusion & Perspectives", "Bilan · Limites · Travaux futurs"),
]

col1 = sections[:4]
col2 = sections[4:]
rh = Inches(0.72)

for i, (num, title, sub) in enumerate(col1):
    ty = Inches(1.1) + i * rh
    b = rect(sl, Inches(0.3), ty, Inches(0.55), rh - Inches(0.06), ACC)
    b.line.fill.background()
    tb(sl, num, Inches(0.3), ty + Inches(0.06), Inches(0.55), rh - Inches(0.1),
       fs=16, bold=True, color=BG, align=PP_ALIGN.CENTER)
    tb(sl, title, Inches(0.97), ty + Inches(0.0), Inches(5.7), Inches(0.4),
       fs=14, bold=True, color=WHITE)
    tb(sl, sub, Inches(0.97), ty + Inches(0.35), Inches(5.7), Inches(0.38),
       fs=10, color=LIGHT, italic=True)

for i, (num, title, sub) in enumerate(col2):
    ty = Inches(1.1) + i * rh
    bx = Inches(7.0)
    b = rect(sl, bx, ty, Inches(0.55), rh - Inches(0.06), ACC)
    b.line.fill.background()
    tb(sl, num, bx, ty + Inches(0.06), Inches(0.55), rh - Inches(0.1),
       fs=16, bold=True, color=BG, align=PP_ALIGN.CENTER)
    tb(sl, title, bx + Inches(0.67), ty, Inches(5.4), Inches(0.4),
       fs=14, bold=True, color=WHITE)
    tb(sl, sub, bx + Inches(0.67), ty + Inches(0.35), Inches(5.4), Inches(0.38),
       fs=10, color=LIGHT, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SECTION: INTRODUCTION
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
section_divider(sl, "01", "Introduction",
    "Contexte · Problématique · Contributions de ce travail")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — CONTEXTE : L'INFODÉMIE
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "Contexte — L'Infodémie COVID-19",
    "OMS, février 2020 : « une surabondance d'informations, exactes ou non, qui rend difficile de trouver des sources fiables »")

# Left: problem description
bullet_box(sl, [
    ("Pandémie de COVID-19 →", "Crise sanitaire mondiale + crise informationnelle parallèle"),
    ("Réseaux sociaux →", "Vecteurs de diffusion rapide de fausses informations : théories du complot, remèdes miracles, désinformation politique"),
    ("Sharma et al. [2] →", "Des millions de publications médicalement erronées par jour en ligne"),
    ("Fact-checking manuel →", "Insuffisant à l'échelle : des dizaines de milliers de posts par heure, impossible à traiter humainement"),
    ("Solution proposée →", "Automatisation par Apprentissage Automatique et NLP"),
],
Inches(0.35), Inches(1.15), Inches(6.8), Inches(5.8), fs=13,
title="Pourquoi ce problème est urgent ?", title_color=ACC)

# Right: stats box
rect(sl, Inches(7.4), Inches(1.15), Inches(5.6), Inches(5.8), DARK, ACC)
stats = [
    ("10 700", "publications COVID-19 annotées"),
    ("~52 %", "de fausses informations détectées"),
    ("3 mois", "de collecte (jan–juil 2020)"),
    ("> 140", "pays touchés par la désinformation"),
    ("1 sur 4", "tweets COVID contenait du contenu trompeur"),
]
tb(sl, "Quelques chiffres clés", Inches(7.6), Inches(1.3), Inches(5.2), Inches(0.45),
   fs=14, bold=True, color=ACC)
for i, (num, desc) in enumerate(stats):
    ty = Inches(1.9) + i * Inches(0.95)
    tb(sl, num, Inches(7.6), ty, Inches(1.6), Inches(0.5),
       fs=30, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    tb(sl, desc, Inches(9.3), ty + Inches(0.05), Inches(3.5), Inches(0.5),
       fs=12, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PROBLÉMATIQUE & CONTRIBUTIONS
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "Problématique & Contributions",
    "Deux questions de recherche complémentaires")

# Two research questions
for i, (q_num, q_title, q_desc, q_color) in enumerate([
    ("RQ1", "Question de Performance",
     "Quelle architecture de modèle ML donne les meilleures performances\npour classer textes COVID-19 comme fake ou real ?",
     ACC),
    ("RQ2", "Question de Sécurité IA",
     "Comment rendre les décisions du modèle transparentes, auditables\net exemptes de biais pour un usage en santé publique ?",
     YELL),
]):
    bx = Inches(0.35) + i * Inches(6.5)
    b = rect(sl, bx, Inches(1.1), Inches(6.2), Inches(1.8), DARK, q_color)
    tb(sl, q_num, bx + Inches(0.15), Inches(1.2), Inches(1.0), Inches(0.5),
       fs=22, bold=True, color=q_color)
    tb(sl, q_title, bx + Inches(1.2), Inches(1.2), Inches(4.8), Inches(0.45),
       fs=14, bold=True, color=WHITE)
    tb(sl, q_desc, bx + Inches(0.15), Inches(1.7), Inches(5.9), Inches(1.1),
       fs=12, color=LIGHT)

# Contributions
tb(sl, "5 Contributions de ce travail", Inches(0.35), Inches(3.2),
   Inches(12.6), Inches(0.45), fs=15, bold=True, color=ACC)

contribs = [
    "Évaluation comparative rigoureuse de 5 modèles ML sur un corpus COVID-19 spécifique",
    "Fine-tuning de DistilBERT avec analyse des courbes d'apprentissage",
    "Analyse d'explicabilité : LIME + visualisation des poids d'attention pour chaque prédiction",
    "Analyse systématique des erreurs → identification de la longueur du texte comme facteur d'échec principal",
    "Pipeline complet packagé en application web interactive (Streamlit + API Flask)",
]
for i, c in enumerate(contribs):
    b = rect(sl, Inches(0.35), Inches(3.75) + i * Inches(0.6),
             Inches(12.6), Inches(0.55), DARK, ACC)
    n = rect(sl, Inches(0.35), Inches(3.75) + i * Inches(0.6),
             Inches(0.45), Inches(0.55), ACC)
    n.line.fill.background()
    tb(sl, str(i+1), Inches(0.35), Inches(3.79) + i * Inches(0.6),
       Inches(0.45), Inches(0.5), fs=13, bold=True, color=BG, align=PP_ALIGN.CENTER)
    tb(sl, c, Inches(0.9), Inches(3.8) + i * Inches(0.6),
       Inches(12.0), Inches(0.5), fs=12, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SECTION: ÉTAT DE L'ART
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
section_divider(sl, "02", "État de l'Art",
    "Évolution de la détection de fake news — 2017 à 2021")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TIMELINE ÉTAT DE L'ART
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "État de l'Art — Évolution 2017–2021",
    "De la détection lexicale aux Transformers contextuels")

# Timeline line
rect(sl, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.06), ACC)

milestones = [
    (Inches(0.6),  "2017", "FakeNews\nChallenge", "TF-IDF +\nSVM / NB\n[Pomerleau]"),
    (Inches(2.9),  "2017", "Marqueurs\nlinguistiques", "Hedging,\nstylométrie\n[Volkova]"),
    (Inches(5.2),  "2018", "BERT", "Pre-training\nbidirectionnel\n[Devlin]"),
    (Inches(7.5),  "2019", "DistilBERT", "Distillation\n40% + léger\n[Sanh]"),
    (Inches(9.8),  "2021", "Constraint\nWorkshop", "COVID-19\nDataset\n[Patwa]"),
    (Inches(11.9), "2021", "LIME/Attention\nsur BERT", "XAI pour\nfake news\n[Popat]"),
]

for i, (mx, year, title, desc) in enumerate(milestones):
    # dot on timeline
    dot = sl.shapes.add_shape(9, mx, Inches(3.35), Inches(0.3), Inches(0.3))
    dot.fill.solid(); dot.fill.fore_color.rgb = ACC
    dot.line.fill.background()

    ty = Inches(1.0) if i % 2 == 0 else Inches(4.1)
    line_top = Inches(1.6) if i % 2 == 0 else Inches(3.7)
    line_bot = Inches(3.35) if i % 2 == 0 else Inches(4.1)

    # vertical connector
    rect(sl, mx + Inches(0.12), line_top, Inches(0.04),
         abs(line_bot - line_top), RGBColor(0x2A, 0x5A, 0x7A))

    rect(sl, mx - Inches(0.1), ty, Inches(1.8), Inches(0.55), DARK, ACC)
    tb(sl, year, mx - Inches(0.1), ty, Inches(1.8), Inches(0.55),
       fs=14, bold=True, color=ACC, align=PP_ALIGN.CENTER)
    tb(sl, title, mx - Inches(0.1), ty + Inches(0.56), Inches(1.9), Inches(0.45),
       fs=11, bold=True, color=WHITE)
    tb(sl, desc, mx - Inches(0.1), ty + Inches(1.0), Inches(1.9), Inches(0.8),
       fs=9.5, color=LIGHT)

# Position of our work
rect(sl, Inches(5.0), Inches(6.0), Inches(3.3), Inches(1.15), DARK, GREEN)
tb(sl, "→ Ce travail (2026)",
   Inches(5.0), Inches(6.0), Inches(3.3), Inches(0.4),
   fs=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
tb(sl, "SVM + TF-IDF vs DistilBERT\n+ LIME + Attention — COVID-19",
   Inches(5.0), Inches(6.42), Inches(3.3), Inches(0.65),
   fs=10, color=LIGHT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — SECTION: DATASET & PRÉTRAITEMENT
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
section_divider(sl, "03", "Dataset & Prétraitement",
    "Constraint COVID-19 · Nettoyage de texte · TF-IDF · Features linguistiques")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — DATASET
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "Dataset — Constraint COVID-19 Fake News",
    "Patwa et al., AAAI Constraint Workshop 2021")

# Left: description
bullet_box(sl, [
    ("Source :", "Publications Twitter collectées entre janvier et juillet 2020"),
    ("Taille :", "10 700 échantillons au total — quasi-équilibré"),
    ("Splits :", "Train : 6 420  |  Validation : 2 140  |  Test : 2 140"),
    ("Labels :", "Binaires — real (info vérifiée) / fake (désinformation confirmée)"),
    ("Annotation :", "Annotateurs formés, sources : OMS, CDC, PubMed, Snopes, PolitiFact"),
    ("Textes :", "Tweets courts — majorité < 50 mots (limite Twitter)"),
],
Inches(0.35), Inches(1.1), Inches(6.0), Inches(5.8), fs=12.5,
title="Caractéristiques du dataset")

# Right: figures
img(sl, os.path.join(FIGURES, '01_class_distribution.png'),
    Inches(6.6), Inches(1.1), Inches(6.5), Inches(2.8))
img(sl, os.path.join(FIGURES, '02_text_length_distribution.png'),
    Inches(6.6), Inches(4.1), Inches(6.5), Inches(3.0))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ANALYSE EXPLORATOIRE
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "Analyse Exploratoire — Vocabulaire par Classe",
    "Les mots révèlent des patterns stylistiques distincts")

img(sl, os.path.join(FIGURES, '03_wordclouds.png'),
    Inches(0.35), Inches(1.1), Inches(7.5), Inches(3.5))
img(sl, os.path.join(FIGURES, '04_top_words.png'),
    Inches(8.0), Inches(1.1), Inches(5.1), Inches(3.5))

tb(sl, "Observations clés", Inches(0.35), Inches(4.75), Inches(12.6), Inches(0.4),
   fs=14, bold=True, color=ACC)

obs = [
    "News RÉELLES : vocabulaire institutionnel — « data », « state », « case », « test », « reported »",
    "FAKE news : vocabulaire alarmiste et politisé — « claim », « cure », « trump », « china », « video »",
    "Les fake news utilisent davantage l'impératif et les titres sensationnalistes",
    "La présence de termes géographiques spécifiques (Chine, USA) est plus marquée dans les fake news",
]
for i, o in enumerate(obs):
    tb(sl, f"• {o}", Inches(0.35), Inches(5.2) + i * Inches(0.48),
       Inches(12.6), Inches(0.44), fs=11.5, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — PRÉTRAITEMENT & FEATURES
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "Prétraitement & Extraction de Features",
    "Pipeline de nettoyage → TF-IDF → Features linguistiques")

# Pipeline diagram
steps = [
    ("1", "Texte brut", "Tweet original"),
    ("2", "Nettoyage", "URLs · @mentions · HTML · ponctuation · minuscules"),
    ("3", "TF-IDF", "15 000 termes · unigrammes + bigrammes · sublinear_tf"),
    ("4", "Features ling.", "7 features numériques : word_count, uppercase_ratio…"),
    ("5", "ColumnTransformer", "Fusion TF-IDF + features numériques (StandardScaler)"),
    ("6", "Matrice finale", "6 420 × 15 007 pour l'entraînement"),
]

arrow_y = Inches(3.5)
step_w  = Inches(2.0)
start_x = Inches(0.3)

for i, (num, title, desc) in enumerate(steps):
    sx = start_x + i * (step_w + Inches(0.15))
    # box
    b = rect(sl, sx, Inches(1.15), step_w, Inches(2.4),
             DARK, ACC if i in [0, 5] else RGBColor(0x00, 0x64, 0x8A))
    tb(sl, num, sx + Inches(0.05), Inches(1.2), Inches(0.4), Inches(0.4),
       fs=18, bold=True, color=ACC)
    tb(sl, title, sx + Inches(0.05), Inches(1.62), step_w - Inches(0.1), Inches(0.42),
       fs=12, bold=True, color=WHITE)
    tb(sl, desc, sx + Inches(0.05), Inches(2.08), step_w - Inches(0.1), Inches(1.4),
       fs=9.5, color=LIGHT)
    # arrow
    if i < len(steps) - 1:
        tb(sl, "→", sx + step_w + Inches(0.02), Inches(2.0), Inches(0.15), Inches(0.5),
           fs=18, bold=True, color=ACC)

# Bottom: 7 linguistic features
tb(sl, "Features Linguistiques Extraites", Inches(0.35), Inches(3.85),
   Inches(12.6), Inches(0.4), fs=14, bold=True, color=ACC)

feats = ["word_count", "char_count", "avg_word_length", "unique_word_ratio",
         "uppercase_ratio", "exclamation_count", "question_count"]
feat_w = Inches(1.75)
for i, f in enumerate(feats):
    b = rect(sl, Inches(0.35) + i * feat_w, Inches(4.35),
             feat_w - Inches(0.08), Inches(0.5), DARK, RGBColor(0x00,0x64,0x8A))
    tb(sl, f, Inches(0.35) + i * feat_w, Inches(4.38),
       feat_w - Inches(0.08), Inches(0.45), fs=10, color=LIGHT, align=PP_ALIGN.CENTER)

tb(sl, "→ Avantage : capture des signaux de surface (majuscules excessives, ponctuation agressive) absents du TF-IDF",
   Inches(0.35), Inches(5.05), Inches(12.6), Inches(0.4), fs=11.5, color=YELL, italic=True)

img(sl, os.path.join(FIGURES, 'feature_distributions.png'),
    Inches(0.35), Inches(5.55), Inches(12.6), Inches(1.75))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — SECTION: ARCHITECTURES
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
section_divider(sl, "04", "Architectures des Modèles",
    "Modèles classiques (SVM · LR · RF · XGB) + Transformer (DistilBERT)")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — MODÈLES CLASSIQUES
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "Modèles Classiques — Architectures & Hyperparamètres",
    "Tous opèrent sur la représentation TF-IDF + features linguistiques")

models_info = [
    ("SVM\n(LinearSVC)", ACC,
     "Séparateur à Vaste Marge\n"
     "• Cherche l'hyperplan de marge maximale\n"
     "• Régularisation L2, C = 1.0\n"
     "• Poids de classe équilibrés\n"
     "• Grid search sur C ∈ {0.01, 0.1, 1, 5, 10}\n"
     "→ Très efficace sur espaces haute dimension"),
    ("Régression\nLogistique", GREEN,
     "Modèle linéaire probabiliste\n"
     "• Régularisation L2, C = 1.0\n"
     "• Solveur lbfgs, 1000 itérations\n"
     "• Produit des probabilités calibrées\n"
     "• AUC-ROC utilisable directement\n"
     "→ Meilleure calibration probabiliste"),
    ("Random\nForest", YELL,
     "Ensemble de 200 arbres\n"
     "• Critère Gini, profondeur max = 20\n"
     "• Bootstrap sampling\n"
     "• Régularisation implicite par moyenne\n"
     "• Résistant au surapprentissage\n"
     "→ Interprétable via feature importance"),
    ("XGBoost", RGBColor(0xFF, 0x8C, 0x00),
     "Gradient Boosting optimisé\n"
     "• 200 estimateurs, depth = 6\n"
     "• Learning rate = 0.1\n"
     "• Régularisation L1 + L2\n"
     "• Correction itérative des erreurs\n"
     "→ Robuste aux features corrélées"),
]

col_w = Inches(3.1)
for i, (name, color, desc) in enumerate(models_info):
    cx = Inches(0.3) + i * (col_w + Inches(0.07))
    # header badge
    b = rect(sl, cx, Inches(1.1), col_w, Inches(0.65), color)
    b.line.fill.background()
    tb(sl, name, cx, Inches(1.12), col_w, Inches(0.62),
       fs=15, bold=True, color=BG, align=PP_ALIGN.CENTER)
    # body
    b2 = rect(sl, cx, Inches(1.78), col_w, Inches(5.3), DARK,
              RGBColor(0x2A, 0x3A, 0x4A))
    tb(sl, desc, cx + Inches(0.1), Inches(1.9), col_w - Inches(0.15), Inches(5.0),
       fs=11, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — DISTILBERT
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "Architecture Transformer — DistilBERT",
    "distilbert-base-uncased — Fine-tuning sur le corpus COVID-19")

# Left: architecture diagram as boxes
arch_items = [
    (ACC,   "Input : Tweet (max 128 tokens)"),
    (DARK,  "[CLS] tok1 tok2 ... [SEP]  → IDs + masque d'attention"),
    (RGBColor(0x00,0x64,0x8A), "DistilBERT Encoder — 6 couches Transformer"),
    (DARK,  "Chaque couche : Multi-Head Self-Attention (12 têtes) + Feed-Forward (3 072 neurones)"),
    (DARK,  "Représentation [CLS] — dim = 768 → résumé sémantique de la séquence"),
    (RGBColor(0x00,0x64,0x8A), "Tête de classification : Dropout(0.3) → Linear(768→256) → ReLU → Dropout(0.3) → Linear(256→2)"),
    (GREEN, "Sortie : Softmax → P(fake) + P(real)"),
]
arch_h = Inches(0.67)
for i, (c, txt) in enumerate(arch_items):
    ty = Inches(1.1) + i * (arch_h + Inches(0.04))
    b = rect(sl, Inches(0.3), ty, Inches(7.2), arch_h, c,
             RGBColor(0x2A, 0x5A, 0x7A) if c == DARK else None)
    if c != DARK and c != GREEN:
        b.line.fill.background()
    tb(sl, txt, Inches(0.4), ty + Inches(0.1), Inches(7.0), arch_h - Inches(0.1),
       fs=11, color=BG if c != DARK else LIGHT, bold=(c != DARK))

# Right: training params
rect(sl, Inches(7.8), Inches(1.1), Inches(5.2), Inches(6.3), DARK, ACC)
tb(sl, "Hyperparamètres d'Entraînement",
   Inches(7.9), Inches(1.2), Inches(5.0), Inches(0.4),
   fs=13, bold=True, color=ACC)

params = [
    ("Optimizer",      "AdamW (weight_decay=0.01)"),
    ("Learning rate",  "2 × 10⁻⁵ (linéaire + warm-up 10 %)"),
    ("Batch size",     "32"),
    ("Epochs",         "3 (early stopping sur val F1)"),
    ("Max tokens",     "128 (padding/truncation)"),
    ("Gradient clip",  "1.0 (stabilisation)"),
    ("Pré-entraîné sur", "Wikipedia + BooksCorpus\n(11 GB de texte)"),
    ("Poids",          "66 M paramètres (vs 110 M BERT)\n→ 40 % plus léger, 60 % plus rapide"),
]
for i, (k, v) in enumerate(params):
    ty = Inches(1.7) + i * Inches(0.58)
    tb(sl, k + " :", Inches(7.9), ty, Inches(1.8), Inches(0.55),
       fs=10.5, bold=True, color=LIGHT)
    tb(sl, v, Inches(9.7), ty, Inches(3.2), Inches(0.55), fs=10.5, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — SECTION: RÉSULTATS
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
section_divider(sl, "05", "Résultats & Discussion",
    "Tableau comparatif · ROC · Matrices de confusion · Feature importance")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — TABLEAU RÉSULTATS + COURBES
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "Résultats — Comparaison Complète des Modèles",
    "Évaluation sur l'ensemble de validation (2 140 échantillons)")

results_table(sl, Inches(0.35), Inches(1.1))

tb(sl, "★ SVM + TF-IDF : meilleur modèle  |  DistilBERT et LR égaux en AUC-ROC (0,976)",
   Inches(0.35), Inches(3.65), Inches(12.6), Inches(0.38),
   fs=11, color=GREEN, italic=True)

img(sl, os.path.join(FIGURES, '07_baseline_comparison.png'),
    Inches(0.35), Inches(4.1), Inches(6.5), Inches(3.2))
img(sl, os.path.join(FIGURES, 'final_results_heatmap.png'),
    Inches(7.0), Inches(4.1), Inches(6.2), Inches(3.2))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — ROC + CONFUSION MATRICES
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "Courbes ROC & Matrices de Confusion",
    "AUC-ROC : Logistic Regression = DistilBERT = 0,976")

img(sl, os.path.join(FIGURES, 'roc_curves.png'),
    Inches(0.35), Inches(1.1), Inches(6.1), Inches(3.4))
img(sl, os.path.join(FIGURES, '06_confusion_matrices_baselines.png'),
    Inches(6.7), Inches(1.1), Inches(6.4), Inches(3.4))

# Key observations
obs_items = [
    "SVM atteint 93,9 % d'accuracy mais ne produit pas de probabilités → pas de courbe ROC directe",
    "LR et DistilBERT ont des AUC-ROC identiques (0,976) → capacité discriminante équivalente",
    "DistilBERT : 72 faux positifs vs 81 faux négatifs — asymétrie importante en santé publique",
    "Les 4 modèles classiques dépassent 90 % d'accuracy → baselines très compétitifs sur ce corpus",
]
tb(sl, "Observations", Inches(0.35), Inches(4.65), Inches(12.6), Inches(0.38),
   fs=13, bold=True, color=ACC)
for i, o in enumerate(obs_items):
    tb(sl, f"• {o}", Inches(0.35), Inches(5.1) + i * Inches(0.5),
       Inches(12.6), Inches(0.46), fs=11.5, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — FEATURE IMPORTANCE
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "Analyse de l'Importance des Features",
    "Coefficients TF-IDF de la Régression Logistique — quels mots distinguent REAL de FAKE ?")

img(sl, os.path.join(FIGURES, '09_feature_importance_lr.png'),
    Inches(0.35), Inches(1.1), Inches(7.8), Inches(5.0))

bullet_box(sl, [
    ("Prédicteurs de REAL :", "'covid19' (hashtags officiels), 'cases', 'testing', 'data', 'reported' → registre institutionnel factuel"),
    ("Prédicteurs de FAKE :", "'trump', 'cure', 'claim', 'video', 'china', 'died' → contenu politisé et alarmiste"),
    ("Observation :", "Le mot 'trump' est un fort prédicteur de fake — politisation marquée du discours COVID-19"),
    ("Attention :", "Ces associations reflètent des biais du corpus, pas des vérités universelles"),
],
Inches(8.3), Inches(1.1), Inches(4.8), Inches(5.0), fs=11.5)

img(sl, os.path.join(FIGURES, '03_wordclouds.png'),
    Inches(0.35), Inches(6.2), Inches(12.6), Inches(1.15))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — SECTION: EXPLICABILITÉ
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
section_divider(sl, "06", "Explicabilité des Modèles (XAI)",
    "LIME · Visualisation de l'attention · Analyse des biais · Considérations éthiques")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — LIME
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "Explicabilité — LIME (Local Interpretable Model-agnostic Explanations)",
    "Ribeiro et al., KDD 2016 — Approximation linéaire locale de la frontière de décision")

# Left: how LIME works
rect(sl, Inches(0.35), Inches(1.1), Inches(5.5), Inches(5.9), DARK, ACC)
tb(sl, "Comment fonctionne LIME ?",
   Inches(0.45), Inches(1.2), Inches(5.3), Inches(0.4), fs=13, bold=True, color=ACC)

lime_steps = [
    "1. Sélectionner un texte à expliquer",
    "2. Générer 500 perturbations locales en masquant des mots aléatoirement",
    "3. Observer les changements de probabilités prédites",
    "4. Ajuster un modèle linéaire sparse sur ces perturbations",
    "5. Afficher les 10 mots les plus influents avec leur poids",
    "",
    "Rouge → pousse vers FAKE",
    "Vert → pousse vers REAL",
]
for i, s in enumerate(lime_steps):
    color = RED if "Rouge" in s else (GREEN if "Vert" in s else LIGHT)
    bold = "→" in s and ("FAKE" in s or "REAL" in s)
    tb(sl, s, Inches(0.45), Inches(1.7) + i * Inches(0.56),
       Inches(5.3), Inches(0.52), fs=11.5, color=color, bold=bold)

# Right: figure
img(sl, os.path.join(FIGURES, 'lime_fake_1.png'),
    Inches(6.1), Inches(1.1), Inches(7.0), Inches(2.9))
img(sl, os.path.join(FIGURES, 'lime_real_1.png'),
    Inches(6.1), Inches(4.15), Inches(7.0), Inches(2.85))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — ATTENTION + BIAIS
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "Attention DistilBERT & Analyse des Biais",
    "Visualisation des poids d'attention + tests de sensibilité aux termes biaisés")

img(sl, os.path.join(FIGURES, '17_attention_visualization.png'),
    Inches(0.35), Inches(1.1), Inches(7.5), Inches(3.0))
img(sl, os.path.join(FIGURES, 'bias_analysis.png'),
    Inches(8.0), Inches(1.1), Inches(5.1), Inches(3.0))

tb(sl, "Résultats de l'analyse des biais", Inches(0.35), Inches(4.25),
   Inches(12.6), Inches(0.38), fs=13, bold=True, color=ACC)

bias_obs = [
    "• Attention (REAL) : distribution diffuse sur mots informatifs (données, annonce, gouvernement)",
    "• Attention (FAKE) : concentration sur termes culturellement chargés (géographiques, ethniques)",
    "• Test de biais : ajouter 'china developed' au texte neutre → augmentation de +4 % de P(fake)",
    "• Implication : le modèle a appris des associations ethniques/géopolitiques du corpus d'entraînement",
    "• ⚠ Risque : déployer ce modèle sans correction pourrait amplifier des discriminations systémiques",
]
for i, b in enumerate(bias_obs):
    color = RED if "⚠" in b else LIGHT
    tb(sl, b, Inches(0.35), Inches(4.7) + i * Inches(0.48),
       Inches(12.6), Inches(0.44), fs=11.5, color=color)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — ANALYSE DES ERREURS
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "Analyse des Erreurs — Facteurs de Difficulté",
    "Quels textes résistent à la classification automatique ?")

img(sl, os.path.join(FIGURES, 'confidence_analysis.png'),
    Inches(0.35), Inches(1.1), Inches(6.3), Inches(3.5))

rect(sl, Inches(6.8), Inches(1.1), Inches(6.2), Inches(3.5), DARK, ACC)
tb(sl, "Distribution de confiance", Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.4),
   fs=13, bold=True, color=ACC)
conf_obs = [
    "Pattern fortement bimodal : la plupart des prédictions sont faites avec très haute confiance (P > 0,90)",
    "Les prédictions incorrectes sont plus uniformément distribuées à des valeurs de confiance intermédiaires",
    "→ L'incertitude du modèle est dans une certaine mesure calibrée : il doute sur les cas vraiment ambigus",
]
for i, o in enumerate(conf_obs):
    tb(sl, f"• {o}", Inches(7.0), Inches(1.7) + i * Inches(0.85),
       Inches(5.8), Inches(0.8), fs=11, color=LIGHT)

tb(sl, "Taux d'erreur vs longueur du texte",
   Inches(0.35), Inches(4.75), Inches(12.6), Inches(0.4),
   fs=13, bold=True, color=ACC)

# Error rate visual
length_bins = [("3–13 mots", "~11 %", RED), ("14–20 mots", "~9 %", YELL),
               ("21–28 mots", "~7 %", YELL), ("29–38 mots", "< 6 %", GREEN)]
bw = Inches(3.0)
for i, (label, rate, color) in enumerate(length_bins):
    bx = Inches(0.35) + i * (bw + Inches(0.07))
    b = rect(sl, bx, Inches(5.25), bw, Inches(0.6), DARK, color)
    tb(sl, label, bx, Inches(5.28), bw, Inches(0.3),
       fs=11, color=LIGHT, align=PP_ALIGN.CENTER)
    tb(sl, rate, bx, Inches(5.6), bw, Inches(0.25),
       fs=13, bold=True, color=color, align=PP_ALIGN.CENTER)

tb(sl, "→ Recommandation : orienter les textes de moins de 15 mots vers une revue humaine",
   Inches(0.35), Inches(6.0), Inches(12.6), Inches(0.38),
   fs=12, color=YELL, italic=True)
tb(sl, "   (taux d'erreur > 10 % sur ces textes — contexte linguistique insuffisant)",
   Inches(0.35), Inches(6.42), Inches(12.6), Inches(0.38),
   fs=11.5, color=LIGHT, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — SECTION: CONCLUSION
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
section_divider(sl, "07", "Conclusion & Perspectives",
    "Bilan des contributions · Limites identifiées · Travaux futurs")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — CONCLUSION
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "Conclusion — Bilan des Contributions",
    "Quatre enseignements principaux de ce travail")

findings = [
    (GREEN, "01", "SVM + TF-IDF = meilleur modèle (93,9 %)",
     "Les modèles classiques bien calibrés restent très compétitifs sur la détection de "
     "désinformation en textes courts à vocabulaire contrôlé. Ne pas les ignorer."),
    (ACC,   "02", "DistilBERT = meilleure calibration probabiliste (AUC 0,976)",
     "La compréhension contextuelle profonde offre une meilleure calibration de probabilité "
     "et une interprétabilité supérieure via l'attention, sans forcément dominer en accuracy."),
    (YELL,  "03", "LIME & Attention révèlent des biais critiques",
     "Les modèles exploitent des marqueurs ethniques et politiques comme prédicteurs de fake. "
     "Ceci constitue un risque de fairness avant tout déploiement."),
    (RED,   "04", "Textes courts = principal mode d'échec",
     "Taux d'erreur > 10 % pour les textes < 13 mots. Recommandation : architecture hybride "
     "avec intervention humaine pour les publications très courtes."),
]

for i, (color, num, title, desc) in enumerate(findings):
    row = i // 2
    col = i % 2
    bx = Inches(0.35) + col * Inches(6.55)
    ty = Inches(1.1) + row * Inches(2.8)
    b = rect(sl, bx, ty, Inches(6.3), Inches(2.6), DARK, color)
    n = rect(sl, bx, ty, Inches(0.6), Inches(2.6), color)
    n.line.fill.background()
    tb(sl, num, bx, ty + Inches(0.8), Inches(0.6), Inches(0.9),
       fs=22, bold=True, color=BG, align=PP_ALIGN.CENTER)
    tb(sl, title, bx + Inches(0.7), ty + Inches(0.15), Inches(5.5), Inches(0.5),
       fs=13, bold=True, color=color)
    tb(sl, desc, bx + Inches(0.7), ty + Inches(0.7), Inches(5.5), Inches(1.7),
       fs=11, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 25 — PERSPECTIVES
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "Perspectives & Travaux Futurs",
    "Quatre axes d'amélioration identifiés")

futures = [
    ("Extension\nMultimodale", ACC,
     "Les fake news sont de plus en plus diffusées avec images et vidéos. "
     "Intégrer des features visuelles (CLIP, ViT) et des métadonnées sociales "
     "(profil de l'auteur, historique de partage) améliorerait significativement les performances."),
    ("Généralisation\nMultilingue", GREEN,
     "L'infodémie COVID-19 était globale — française, arabe, mandarine. "
     "Fine-tuner des modèles multilingues (mBERT, XLM-R) sur des corpus "
     "multilingues (CLEF CheckThat!) élargirait l'applicabilité du système."),
    ("Robustesse\nTemporelle", YELL,
     "La désinformation évolue dans le temps (des origines du virus à l'hésitation "
     "vaccinale). Des stratégies d'apprentissage continu ou de réentraînement "
     "périodique maintiendraient la performance face à ces changements."),
    ("Mitigation\ndes Biais", RED,
     "L'analyse d'explicabilité révèle des biais démographiques. "
     "L'augmentation de données contrefactuelles et le débiaisage adversarial "
     "pourraient réduire la dépendance aux identificateurs géographiques/ethniques."),
]

for i, (name, color, desc) in enumerate(futures):
    col = i % 2
    row = i // 2
    bx = Inches(0.35) + col * Inches(6.55)
    ty = Inches(1.1) + row * Inches(2.8)
    b = rect(sl, bx, ty, Inches(6.3), Inches(2.6), DARK, color)
    tb(sl, name, bx + Inches(0.15), ty + Inches(0.15),
       Inches(6.0), Inches(0.65), fs=15, bold=True, color=color)
    hline_s = rect(sl, bx + Inches(0.15), ty + Inches(0.82),
                   Inches(5.9), Pt(1.5), color)
    hline_s.line.fill.background()
    tb(sl, desc, bx + Inches(0.15), ty + Inches(0.92),
       Inches(6.0), Inches(1.55), fs=11, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 26 — RÉFÉRENCES
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "Références Bibliographiques")

refs_left = [
    "[1] OMS, Situation Report, 2020.",
    "[2] Sharma et al., ACM TIST 2019.",
    "[3] Pérez-Rosas et al., COLING 2018.",
    "[4] Devlin et al. BERT, NAACL 2019.",
    "[5] Sanh et al. DistilBERT, arXiv 2019.",
    "[6] Ribeiro et al. LIME, KDD 2016.",
    "[7] Lundberg & Lee SHAP, NeurIPS 2017.",
    "[8] Pomerleau, FakeNewsChallenge 2017.",
    "[9] Rashkin et al., EMNLP 2017.",
]
refs_right = [
    "[10] Volkova et al., ACL 2017.",
    "[11] Kula et al., ICCS 2020.",
    "[12] Goldani et al., App. Soft Comp. 2021.",
    "[13] Patwa et al. Constraint, AAAI 2021.",
    "[14] Popat et al., WWW 2018.",
    "[15] Jain & Wallace, NAACL 2019.",
    "[16] Joachims, ECML 1998.",
    "[17] Nakov et al., CLEF 2021.",
]

for i, r in enumerate(refs_left):
    tb(sl, r, Inches(0.4), Inches(1.1) + i * Inches(0.62),
       Inches(6.3), Inches(0.58), fs=11, color=LIGHT)
for i, r in enumerate(refs_right):
    tb(sl, r, Inches(7.0), Inches(1.1) + i * Inches(0.62),
       Inches(6.1), Inches(0.58), fs=11, color=LIGHT)

hline(sl, Inches(6.65))
tb(sl, "Pipeline complet, modèles entraînés et application web disponibles sur GitHub : github.com/Youssefmab/fake-news-detection-project",
   Inches(0.35), Inches(6.75), Inches(12.6), Inches(0.4),
   fs=11, color=ACC, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved : {OUT}")
print(f"Slides : {len(prs.slides)}")
